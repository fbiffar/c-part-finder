from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from PIL import Image

def create_pptx_with_annotations(image_path, annotations, output_path="annotated_presentation.pptx"):
    """
    Generate a PowerPoint slide with the given image and annotations.
    
    Parameters:
    - image_path: Path to the image to be inserted.
    - annotations: List of tuples containing annotation details [(coords, part_name), ...].
    - output_path: Path where the generated PowerPoint file will be saved.
    """
    # Load the image to get its dimensions
    img = Image.open(image_path)
    img_width, img_height = img.size

    # Create PowerPoint presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Calculate the desired image width (50% of slide width) and maintain aspect ratio
    image_width_on_slide = slide_width * 0.5
    image_aspect_ratio = img_height / img_width
    image_height_on_slide = image_width_on_slide * image_aspect_ratio

    # Center the image on the slide
    image_left = (slide_width - image_width_on_slide) / 2  # Center horizontally
    image_top = (slide_height - image_height_on_slide) / 2  # Center vertically
    slide_image = slide.shapes.add_picture(image_path, image_left, image_top, width=image_width_on_slide, height=image_height_on_slide)

    # Calculate scaling factors for annotation placement
    scale_x = image_width_on_slide / img_width
    scale_y = image_height_on_slide / img_height

    # Calculate free space for labels (outside the image)
    free_space_left = 0 if image_left > slide_width / 2 else image_left + image_width_on_slide
    free_space_width = abs(slide_width - image_width_on_slide) / 2

    # Center labels vertically in the free space
    label_top_start = (slide_height - len(annotations) * Inches(0.7)) / 2  # Dynamically position labels

    # Add annotations
    for idx, (coords, part_name) in enumerate(annotations):
        # Unpack coordinates
        left, top, right, bottom = coords

        # Scale coordinates to match PowerPoint dimensions
        left_scaled = left * scale_x + image_left
        top_scaled = top * scale_y + image_top
        right_scaled = right * scale_x + image_left
        bottom_scaled = bottom * scale_y + image_top

        # Draw the box around the part
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left_scaled, top_scaled,
            right_scaled - left_scaled, bottom_scaled - top_scaled
        )
        shape.line.color.rgb = RGBColor(255, 0, 0)  # Red border
        shape.fill.background()  # Transparent fill

        # Position label dynamically in free space
        label_left = free_space_left + Inches(0.2)  # Add margin
        label_top = label_top_start + idx * Inches(0.7)  # Stack labels vertically

        # Add an arrow from the rectangle to the label
        arrow = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, right_scaled, (top_scaled + bottom_scaled) / 2,
            label_left, label_top + Inches(0.25)
        )
        arrow.line.color.rgb = RGBColor(0, 0, 255)  # Blue arrow

        # Add the label with the part name
        label = slide.shapes.add_textbox(label_left, label_top, free_space_width - Inches(0.4), Inches(0.5))
        text_frame = label.text_frame
        text_frame.text = part_name
        text_frame.paragraphs[0].font.size = Pt(12)  # Use Pt for font size
        text_frame.paragraphs[0].font.bold = True

    # Save the PowerPoint presentation
    prs.save(output_path)
    print(f"PowerPoint saved to {output_path}")


# Example Usage
image_path = "c-part-finder/images/conveyor_machine.png"
annotations = [
    # Example annotations [(left, top, right, bottom), "part name"]
    [(100, 150, 200, 250), "Screw"],
    [(300, 350, 400, 450), "Bolt"],
    [(150,50,200,100), "Hinge"],
]
output_path = "annotated_presentation_dynamic_labels_centered.pptx"

create_pptx_with_annotations(image_path, annotations, output_path)
