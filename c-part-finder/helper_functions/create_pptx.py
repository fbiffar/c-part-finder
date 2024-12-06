from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from PIL import Image
import numpy as np
import io
from helper_functions.json_handler import *

from pptx.oxml.ns import qn
from pptx.oxml import parse_xml


def create_pptx_with_annotations(img, annotations, output_path="annotated_presentation.pptx"):
    """
    Generate a PowerPoint slide with the given image and annotations.

    Parameters:
    - img: PIL.Image instance to be inserted.
    - annotations: List of tuples containing annotation details [(coords, category_name, subcategory_name, url), ...].
    - output_path: Path where the generated PowerPoint file will be saved.
    """
    json_path = "c-part-finder/structured_categories/restructured_categories.json"

    # Convert the PIL Image to a file-like object (BytesIO)
    img_bytes = io.BytesIO()
    img.save(img_bytes, format="PNG")
    img_bytes.seek(0)  # Reset the file pointer to the beginning

    img_width, img_height = img.size

    # Create PowerPoint presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Calculate the desired image width (50% of slide width) and maintain aspect ratio
    image_width_on_slide = slide_width * 0.5
    image_aspect_ratio = img_height / img_width
    image_height_on_slide = image_width_on_slide * image_aspect_ratio

    # Center the image on the slide and add a black frame
    image_left = (slide_width - image_width_on_slide) / 2  # Center horizontally
    image_top = (slide_height - image_height_on_slide) / 2  # Center vertically
    image_shape = slide.shapes.add_picture(
        img_bytes, image_left, image_top, width=image_width_on_slide, height=image_height_on_slide
    )
    image_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
    image_shape.line.width = Pt(3)  # Thickness of the border

    # Calculate scaling factors for annotation placement
    scale_x = image_width_on_slide / img_width
    scale_y = image_height_on_slide / img_height

    # Free space dimensions for labels and images
    free_space_width = Inches(1.8)  # Fixed width for labels and images
    free_space_height = slide_height - Inches(1.5)  # Vertical space for labels

    # Sort annotations by horizontal center of their marked region
    sorted_annotations = sorted(
        annotations,
        key=lambda ann: ((ann[0][0] + ann[0][2]) / 2) * scale_x + image_left
    )

    # Split sorted annotations into left and right groups
    mid_index = len(sorted_annotations) // 2
    left_annotations = sorted_annotations[:mid_index]
    right_annotations = sorted_annotations[mid_index:]

    # Assign labels dynamically based on their sorted positions
    label_positions = {"left": left_annotations, "right": right_annotations}

    print(label_positions)

    # Helper to add labels and arrows
    def add_label_and_arrow(side, start_top, annotations):
        if not annotations:
            return  # Skip if no annotations for this side

        current_top = start_top
        label_spacing = (free_space_height - start_top) / len(annotations)

        for coords, category_name, subcategory_name, url in annotations:
            # Extract part details from JSON
            part_info = extract_category_details(json_path, category_name, subcategory_name)
            if part_info:
                path_image = part_info[3]
                img_part = Image.open(path_image)

            # Prepare label and calculate its position
            label_left = (
                image_left + image_width_on_slide + Inches(0.2)
                if side == "right"
                else image_left - free_space_width - Inches(0.2)
            )
            label_center_x = label_left + free_space_width / 2
            label_height = Inches(0.7)  # Allow for 2 lines of text
            label_top = current_top
            label = slide.shapes.add_textbox(label_left, label_top, free_space_width, label_height)
            text_frame = label.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = f"{subcategory_name}"
            p.font.size = Pt(12)
            p.font.bold = True
            text_frame.paragraphs[0].alignment = 1  # Center alignment

            # Adjust arrow direction dynamically
            left, top, right, bottom = coords
            left_scaled = left * scale_x + image_left
            right_scaled = right * scale_x + image_left
            top_scaled = top * scale_y + image_top
            bottom_scaled = bottom * scale_y + image_top
            region_center_y = (top_scaled + bottom_scaled) / 2
            arrow_start_x = right_scaled if side == "right" else left_scaled
            arrow_start_y = region_center_y
            arrow_end_x = label_left if side == "right" else label_left + free_space_width
            arrow_end_y = label_top + label_height 

            # Add arrow
            arrow = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                arrow_start_x,
                arrow_start_y,
                arrow_end_x,
                arrow_end_y,
            )
            arrow.line.color.rgb = RGBColor(0, 0, 255)  # Blue arrow

            # Draw the box around the part
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left_scaled, top_scaled,
                right_scaled - left_scaled, bottom_scaled - top_scaled
            )
            shape.line.color.rgb = RGBColor(255, 0, 0)  # Red border
            shape.line.width = Pt(3)  # Thickness of the border

            shape.fill.background()  # Transparent fill

             # Add the subcategory image below the label
            if path_image:
                img_part = Image.open(path_image)
                img_bytes_part = io.BytesIO()
                img_part.save(img_bytes_part, format="PNG")
                img_bytes_part.seek(0)
            
                # Calculate the size for the small image
                part_image_width = Inches(1.5)
                part_aspect_ratio = img_part.height / img_part.width
                part_image_height = part_image_width * part_aspect_ratio

                img_shape = slide.shapes.add_picture(
                    img_bytes_part,
                    label_left,
                    label_top + label_height + Inches(0.1),  # Position below the label
                    width=part_image_width,
                    height=part_image_height
                )
                img_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
                img_shape.line.width = Pt(3)  # Thickness of the border
            # Increment top position for the next label
            current_top += label_spacing

    # Add left and right labels and arrows
    add_label_and_arrow("left", Inches(1.0), label_positions["left"])
    add_label_and_arrow("right", Inches(1.0), label_positions["right"])

    # Save the PowerPoint presentation
    prs.save(output_path)
    print(f"PowerPoint saved to {output_path}")



# annotation = [((112, 215, 187, 301), 'Flüssigkeitsmanagement', 'Hydraulische Komponenten', 'https://www.bossard.com/eshop/ch-de/produkte/funktionselemente/fluessigkeitsmanagement/hydraulische-komponenten/c/03.300.200/'), ((196, 38, 251, 86), 'Norm- und Standard Verbindungselemente', 'Schrauben', 'https://www.bossard.com/eshop/ch-de/produkte/verbindungstechnik/norm-und-standard-verbindungselemente/schrauben/c/01.100.100/')]
# img_path = "c-part-finder/images/conveyor_machine.png"
# img = Image.open(img_path)
# create_pptx_with_annotations(img, annotation, output_path="annotated_presentation.pptx")